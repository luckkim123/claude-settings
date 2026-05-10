#!/usr/bin/env python3
"""
verify_pptx.py — PowerPoint .pptx 다층 무결성 검사 (V2.1~V2.5)

5개 항목을 모두 통과해야 OK. 하나라도 실패하면 exit 1 + 진단 출력.

사용:
    python3 verify_pptx.py <path/to/deck.pptx>
    python3 verify_pptx.py <path> --json   # 머신 판독용 JSON

설계 근거: docs/plans/2026-05-11-ppt-edit-analyze-design.md §3.4
v20 사고(orphan slideMaster)로부터 도출. soffice/python-pptx만 통과해도
PowerPoint 복구 다이얼로그가 뜨는 사례를 차단.
"""
from __future__ import annotations
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import dataclass, field, asdict


@dataclass
class CheckResult:
    name: str
    passed: bool
    detail: str = ""
    skipped: bool = False


@dataclass
class VerifyReport:
    pptx_path: str
    checks: list[CheckResult] = field(default_factory=list)

    @property
    def all_passed(self) -> bool:
        return all(c.passed or c.skipped for c in self.checks)

    @property
    def has_skipped(self) -> bool:
        return any(c.skipped for c in self.checks)


def _check_v21_zip(pptx_path: str) -> CheckResult:
    """V2.1 — zip 무결성 (python -m zipfile -t)"""
    try:
        with zipfile.ZipFile(pptx_path) as z:
            bad = z.testzip()
        if bad is None:
            return CheckResult("V2.1 zip 무결성", True, "zip CRC OK")
        return CheckResult("V2.1 zip 무결성", False, f"손상 entry: {bad}")
    except Exception as e:
        return CheckResult("V2.1 zip 무결성", False, f"zip 오픈 실패: {e}")


def _check_v22_python_pptx(pptx_path: str) -> CheckResult:
    """V2.2 — python-pptx Presentation() open"""
    try:
        from pptx import Presentation
    except ImportError:
        return CheckResult(
            "V2.2 python-pptx open",
            False,
            "python-pptx 미설치 (pip3 install --user python-pptx)",
            skipped=True,
        )
    try:
        p = Presentation(pptx_path)
        n = len(p.slides)
        return CheckResult("V2.2 python-pptx open", True, f"slides={n}")
    except Exception as e:
        return CheckResult("V2.2 python-pptx open", False, f"파싱 실패: {e}")


def _check_v23_soffice(pptx_path: str) -> CheckResult:
    """V2.3 — soffice PDF 변환"""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice is None:
        return CheckResult(
            "V2.3 soffice PDF 변환",
            False,
            "soffice 미설치 (brew install --cask libreoffice)",
            skipped=True,
        )
    with tempfile.TemporaryDirectory() as tmp:
        try:
            r = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf",
                 "--outdir", tmp, pptx_path],
                capture_output=True, text=True, timeout=120,
            )
            pdfs = [f for f in os.listdir(tmp) if f.endswith(".pdf")]
            if r.returncode == 0 and pdfs:
                return CheckResult("V2.3 soffice PDF 변환", True, f"PDF={pdfs[0]}")
            return CheckResult(
                "V2.3 soffice PDF 변환", False,
                f"변환 실패 rc={r.returncode}: {r.stderr.strip()[-200:]}"
            )
        except subprocess.TimeoutExpired:
            return CheckResult("V2.3 soffice PDF 변환", False, "타임아웃 (120s)")
        except Exception as e:
            return CheckResult("V2.3 soffice PDF 변환", False, f"실행 실패: {e}")


def _check_v24_ooxml(pptx_path: str) -> CheckResult:
    """V2.4 — OOXML 직접 검사 (sldIdLst, rels, Content_Types)"""
    issues: list[str] = []
    try:
        with zipfile.ZipFile(pptx_path) as z:
            names = set(z.namelist())

            # (a) sldIdLst id ≥256, unique
            try:
                pres = z.read("ppt/presentation.xml").decode("utf-8")
            except KeyError:
                return CheckResult(
                    "V2.4 OOXML 직접 검사", False,
                    "ppt/presentation.xml 없음"
                )
            sldids = re.findall(
                r'<p:sldId\s+id="(\d+)"\s+r:id="([^"]+)"', pres
            )
            ids = [int(i) for i, _ in sldids]
            if not ids:
                issues.append("sldIdLst 비어 있음")
            else:
                if min(ids) < 256:
                    issues.append(f"sldId<256 발견: {min(ids)}")
                dups = [i for i in set(ids) if ids.count(i) > 1]
                if dups:
                    issues.append(f"sldId 중복: {dups}")

            # (b) presentation.xml.rels와 sldIdLst 1:1
            try:
                pres_rels = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
            except KeyError:
                return CheckResult(
                    "V2.4 OOXML 직접 검사", False,
                    "ppt/_rels/presentation.xml.rels 없음"
                )
            slide_rels = re.findall(
                r'Id="([^"]+)"\s+Type="[^"]*relationships/slide"\s+Target="([^"]+)"',
                pres_rels,
            )
            rid_to_target = {r: t for r, t in slide_rels}
            for sid, rid in sldids:
                target = rid_to_target.get(rid)
                if target is None:
                    issues.append(f"sldId {sid} rid={rid}: rel 없음")
                    continue
                full = f"ppt/{target}"
                if full not in names:
                    issues.append(f"sldId {sid}: 파일 {full} 없음")

            # (c) slide files vs pres.rels (양방향)
            slide_files = {n for n in names if n.startswith("ppt/slides/slide")
                           and n.endswith(".xml") and "/_rels/" not in n}
            referenced = {f"ppt/{t}" for _, t in slide_rels}
            orphan_files = slide_files - referenced
            if orphan_files:
                issues.append(f"slide 파일은 있는데 pres.rels에 없음: {sorted(orphan_files)}")
            phantom = referenced - slide_files
            if phantom:
                issues.append(f"pres.rels에 있는데 파일 없음: {sorted(phantom)}")

            # (d) slide.xml.rels의 모든 Target 존재 (dangling rels)
            for n in names:
                if not (n.startswith("ppt/slides/_rels/") and n.endswith(".rels")):
                    continue
                body = z.read(n).decode("utf-8")
                for t in re.findall(r'Target="([^"]+)"', body):
                    if t.startswith("http"):
                        continue
                    if t.startswith("../"):
                        path = "ppt/" + t[3:]
                    else:
                        path = "ppt/slides/" + t
                    path = os.path.normpath(path).replace(os.sep, "/")
                    if path not in names:
                        issues.append(f"{n} dangling: {t}")

            # (e) [Content_Types].xml: 모든 .xml part가 Override 또는 Default 처리
            try:
                ct = z.read("[Content_Types].xml").decode("utf-8")
            except KeyError:
                issues.append("[Content_Types].xml 없음")
                return CheckResult(
                    "V2.4 OOXML 직접 검사",
                    not issues, "; ".join(issues) if issues else "OK"
                )
            defaults = dict(re.findall(
                r'<Default\s+Extension="([^"]+)"\s+ContentType="([^"]+)"', ct
            ))
            overrides = set(re.findall(r'PartName="([^"]+)"', ct))
            for n in names:
                if n.endswith("/"):
                    continue
                if n.startswith("[Content_Types]"):
                    continue
                if n.endswith(".xml") or n.endswith(".rels"):
                    ext = n.rsplit(".", 1)[-1]
                    pn = "/" + n
                    if pn not in overrides and ext not in defaults:
                        issues.append(f"Content_Types 누락: {pn}")
    except Exception as e:
        return CheckResult("V2.4 OOXML 직접 검사", False, f"검사 실패: {e}")

    if issues:
        return CheckResult(
            "V2.4 OOXML 직접 검사", False,
            f"{len(issues)}개 위반 — " + "; ".join(issues[:5])
            + ("…" if len(issues) > 5 else "")
        )
    return CheckResult("V2.4 OOXML 직접 검사", True, "sldIdLst/rels/Content_Types OK")


def _check_v25_orphan_master(pptx_path: str) -> CheckResult:
    """V2.5 — Orphan slideMaster/slideLayout 검사 (v20 사고 재발 차단)

    모든 slideMaster가 ≥1 slide에서 사용되어야 함.
    사용 안 되는 master/layout 발견 시 PowerPoint 복구 다이얼로그 트리거.
    """
    try:
        with zipfile.ZipFile(pptx_path) as z:
            names = z.namelist()
            # master, layout 파일
            master_files = sorted(
                n for n in names
                if n.startswith("ppt/slideMasters/slideMaster")
                and n.endswith(".xml") and "/_rels/" not in n
            )
            layout_files = sorted(
                n for n in names
                if n.startswith("ppt/slideLayouts/slideLayout")
                and n.endswith(".xml") and "/_rels/" not in n
            )

            # slide → layout → master 추적
            master_use: dict[str, int] = {m: 0 for m in master_files}
            layout_use: dict[str, int] = {L: 0 for L in layout_files}

            slide_rels = sorted(
                n for n in names
                if n.startswith("ppt/slides/_rels/slide") and n.endswith(".rels")
            )
            for sr in slide_rels:
                body = z.read(sr).decode("utf-8")
                m = re.search(r'Target="\.\./slideLayouts/(slideLayout\d+\.xml)"', body)
                if not m:
                    continue
                layout_path = f"ppt/slideLayouts/{m.group(1)}"
                if layout_path in layout_use:
                    layout_use[layout_path] += 1
                # 그 layout의 master
                lr_path = f"ppt/slideLayouts/_rels/{m.group(1)}.rels"
                if lr_path in names:
                    lr_body = z.read(lr_path).decode("utf-8")
                    mm = re.search(
                        r'Target="\.\./slideMasters/(slideMaster\d+\.xml)"', lr_body
                    )
                    if mm:
                        master_path = f"ppt/slideMasters/{mm.group(1)}"
                        if master_path in master_use:
                            master_use[master_path] += 1

            orphan_masters = [m for m, c in master_use.items() if c == 0]
            orphan_layouts = [L for L, c in layout_use.items() if c == 0]

            issues = []
            if orphan_masters:
                issues.append(
                    f"orphan slideMaster {len(orphan_masters)}개: "
                    + ", ".join(os.path.basename(m) for m in orphan_masters)
                )
            if orphan_layouts:
                # 정상 .pptx도 안 쓰는 layout은 많이 가지므로 경고만 (실패 X)
                # — 원인이 PowerPoint 다이얼로그를 띄우는 건 master 단위
                pass

            if orphan_masters:
                return CheckResult(
                    "V2.5 Orphan master 검사", False,
                    f"{issues[0]} → PowerPoint 복구 다이얼로그 원인. "
                    f"presentation.xml.rels + sldMasterIdLst + Content_Types Override "
                    f"동시 정리 필요."
                )
            return CheckResult(
                "V2.5 Orphan master 검사", True,
                f"slideMaster {len(master_files)}개 모두 사용됨 "
                f"(layout {len(layout_files)}개 중 {sum(1 for c in layout_use.values() if c>0)}개 사용)"
            )
    except Exception as e:
        return CheckResult("V2.5 Orphan master 검사", False, f"검사 실패: {e}")


def verify(pptx_path: str) -> VerifyReport:
    if not os.path.exists(pptx_path):
        rep = VerifyReport(pptx_path)
        rep.checks.append(CheckResult("파일 존재", False, "파일 없음"))
        return rep
    rep = VerifyReport(pptx_path)
    rep.checks.append(_check_v21_zip(pptx_path))
    rep.checks.append(_check_v22_python_pptx(pptx_path))
    rep.checks.append(_check_v23_soffice(pptx_path))
    rep.checks.append(_check_v24_ooxml(pptx_path))
    rep.checks.append(_check_v25_orphan_master(pptx_path))
    return rep


def _print_human(rep: VerifyReport) -> None:
    print(f"\n📋 verify_pptx — {rep.pptx_path}\n")
    for c in rep.checks:
        if c.skipped:
            mark = "⊘ SKIP"
        elif c.passed:
            mark = "✓ PASS"
        else:
            mark = "✗ FAIL"
        print(f"  {mark}  {c.name}")
        if c.detail:
            print(f"          {c.detail}")
    print()
    if rep.all_passed:
        if rep.has_skipped:
            print("⚠️  통과 (일부 SKIP — 누락 도구 설치 후 재검증 권장)")
        else:
            print("✅ 통과 (5/5)")
    else:
        n_fail = sum(1 for c in rep.checks if not c.passed and not c.skipped)
        print(f"❌ 실패 ({n_fail}개 항목 위반)")


def main() -> int:
    args = sys.argv[1:]
    if not args or args[0] in ("-h", "--help"):
        print(__doc__)
        return 0
    pptx_path = args[0]
    as_json = "--json" in args
    rep = verify(pptx_path)
    if as_json:
        out = {
            "pptx_path": rep.pptx_path,
            "all_passed": rep.all_passed,
            "has_skipped": rep.has_skipped,
            "checks": [asdict(c) for c in rep.checks],
        }
        print(json.dumps(out, ensure_ascii=False, indent=2))
    else:
        _print_human(rep)
    return 0 if rep.all_passed else 1


if __name__ == "__main__":
    sys.exit(main())
