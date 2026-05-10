#!/usr/bin/env python3
"""
extract_spec.py — .pptx 1개에서 4개 산출물 추출 (G4 hybrid 기본, --detail full로 G3)

산출물 (output/<slug>/):
  A. style_spec.yaml       — design tokens (typography/colors/spacing) + layout별 spec
  B. layout_catalog.md     — 슬라이드별 layout 패턴 분류 + 빈도
  C. narrative_outline.md  — 슬라이드 제목 + bullet 트리
  D. asset_manifest.yaml   — 사용된 이미지/차트/폰트 목록
  (옵션) E. slides_full_dump.yaml — --detail full 시 슬라이드별 완전 dump

사용:
    python3 extract_spec.py <deck.pptx> [--out OUTDIR] [--detail full]

설계 근거: docs/plans/2026-05-11-ppt-edit-analyze-design.md §4
"""
from __future__ import annotations
import argparse
import os
import sys
from collections import Counter, defaultdict
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    sys.stderr.write(
        "ERROR: python-pptx 미설치. `pip3 install --user python-pptx`\n"
    )
    sys.exit(1)


def _emu_to_in(v) -> float:
    if v is None:
        return 0.0
    return round(v / 914400, 3)


def _color_hex(color) -> str | None:
    """python-pptx ColorFormat → hex (or None if not solid RGB)."""
    if color is None:
        return None
    try:
        if color.type is None:
            return None
        rgb = color.rgb
        if rgb is None:
            return None
        return f"#{str(rgb).upper()}"
    except Exception:
        return None


def _font_info(run) -> dict:
    f = run.font
    info: dict = {
        "family": f.name,
        "size_pt": f.size.pt if f.size else None,
        "bold": f.bold,
        "italic": f.italic,
        "color_hex": None,
    }
    try:
        info["color_hex"] = _color_hex(f.color)
    except Exception:
        pass
    return info


def _shape_info(sh, with_runs: bool = True) -> dict:
    info: dict = {
        "shape_id": sh.shape_id,
        "name": sh.name,
        "type": str(sh.shape_type),
        "pos_in": [_emu_to_in(sh.left), _emu_to_in(sh.top)],
        "size_in": [_emu_to_in(sh.width), _emu_to_in(sh.height)],
    }
    if sh.has_text_frame:
        runs = []
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if not run.text:
                    continue
                runs.append({"text": run.text, "font": _font_info(run)})
        info["runs"] = runs if with_runs else f"{len(runs)} runs"
        info["text_summary"] = " ".join(
            r["text"] for r in runs
        )[:80] if runs else ""
    return info


def _classify_layout(slide) -> str:
    """슬라이드의 layout 패턴 자동 분류 (휴리스틱)."""
    n_text = sum(1 for s in slide.shapes if s.has_text_frame)
    n_pic = sum(1 for s in slide.shapes if s.shape_type == 13)  # PICTURE
    n_total = len(slide.shapes)
    if n_total <= 2 and n_text >= 1:
        return "title_only"
    if n_pic == 0 and n_text >= 2:
        # 좌우 분할 추정
        xs = sorted(_emu_to_in(s.left) for s in slide.shapes if s.has_text_frame)
        if len(xs) >= 2 and xs[-1] - xs[0] > 4.5:
            return "two_column_text"
        return "title_with_content"
    if n_pic >= 1:
        return "image_with_caption" if n_text >= 1 else "image_only"
    return "other"


def extract(pptx_path: str, outdir: Path, detail_full: bool = False) -> None:
    p = Presentation(pptx_path)
    outdir.mkdir(parents=True, exist_ok=True)

    # === 집계 ===
    layout_counter: Counter = Counter()
    layout_examples: dict[str, list[int]] = defaultdict(list)
    layout_specs: dict[str, dict] = {}
    font_families: Counter = Counter()
    font_sizes: Counter = Counter()
    colors_used: Counter = Counter()
    images: list[dict] = []

    title_runs: list[dict] = []
    body_runs: list[dict] = []
    slides_dump: list[dict] = []

    for idx, slide in enumerate(p.slides, start=1):
        layout_name = _classify_layout(slide)
        layout_counter[layout_name] += 1
        layout_examples[layout_name].append(idx)

        slide_info: dict = {"slide_no": idx, "layout_used": layout_name, "shapes": []}

        for sh in slide.shapes:
            sh_info = _shape_info(sh, with_runs=True)
            slide_info["shapes"].append(sh_info)

            # 폰트/컬러 집계
            if "runs" in sh_info and isinstance(sh_info["runs"], list):
                for r in sh_info["runs"]:
                    f = r["font"]
                    if f["family"]:
                        font_families[f["family"]] += 1
                    if f["size_pt"]:
                        font_sizes[round(f["size_pt"])] += 1
                    if f["color_hex"]:
                        colors_used[f["color_hex"]] += 1
                    # title vs body 추정 (위치 기반)
                    top_in = sh_info["pos_in"][1]
                    if top_in < 1.5:
                        title_runs.append(r)
                    else:
                        body_runs.append(r)

            # 이미지 (Picture shape만 .image 속성 보유)
            if sh.shape_type == 13:
                try:
                    img = getattr(sh, "image", None)
                    if img is not None:
                        images.append({
                            "slide_no": idx,
                            "shape_name": sh.name,
                            "ext": img.ext,
                            "size_bytes": len(img.blob),
                            "pos_in": sh_info["pos_in"],
                            "size_in": sh_info["size_in"],
                        })
                except Exception:
                    pass

        slides_dump.append(slide_info)

        # layout spec 첫 등장 슬라이드를 대표로
        if layout_name not in layout_specs:
            layout_specs[layout_name] = {
                "first_slide": idx,
                "shapes": [
                    {k: v for k, v in s.items() if k != "runs"}
                    | (
                        {"font_summary": _summarize_runs(s.get("runs", []))}
                        if isinstance(s.get("runs"), list) else {}
                    )
                    for s in slide_info["shapes"]
                ],
            }

    # === 산출물 A: style_spec.yaml ===
    title_font = _summarize_runs(title_runs, top_n=1)
    body_font = _summarize_runs(body_runs, top_n=1)
    palette = [c for c, _ in colors_used.most_common(8)]

    style_spec = {
        "deck": {
            "source": os.path.basename(pptx_path),
            "n_slides": len(p.slides),
            "slide_size_in": [
                _emu_to_in(p.slide_width),
                _emu_to_in(p.slide_height),
            ],
            "design_tokens": {
                "typography": {
                    "title_default": title_font,
                    "body_default": body_font,
                    "all_families_used": [
                        {"family": f, "count": c}
                        for f, c in font_families.most_common()
                    ],
                    "all_sizes_used_pt": [
                        {"size_pt": s, "count": c}
                        for s, c in sorted(font_sizes.items())
                    ],
                },
                "palette_hex": palette,
            },
        },
        "layouts": {
            name: {
                "count": layout_counter[name],
                "example_slides": layout_examples[name][:5],
                "spec": layout_specs[name],
            }
            for name in layout_counter
        },
    }
    _write_yaml(outdir / "style_spec.yaml", style_spec)

    # === 산출물 B: layout_catalog.md ===
    md = ["# Layout Catalog\n", f"Source: `{os.path.basename(pptx_path)}` ({len(p.slides)} slides)\n"]
    md.append("\n| Layout | Count | Example Slides |\n|:--|:--:|:--|\n")
    for name, count in layout_counter.most_common():
        examples = ", ".join(str(s) for s in layout_examples[name][:8])
        md.append(f"| `{name}` | {count} | {examples} |\n")
    (outdir / "layout_catalog.md").write_text("".join(md), encoding="utf-8")

    # === 산출물 C: narrative_outline.md ===
    md = ["# Narrative Outline\n", f"Source: `{os.path.basename(pptx_path)}`\n\n"]
    for s in slides_dump:
        md.append(f"## Slide {s['slide_no']} — `{s['layout_used']}`\n\n")
        for sh in s["shapes"]:
            txt = sh.get("text_summary", "")
            if txt:
                md.append(f"- {txt}\n")
        md.append("\n")
    (outdir / "narrative_outline.md").write_text("".join(md), encoding="utf-8")

    # === 산출물 D: asset_manifest.yaml ===
    manifest = {
        "source": os.path.basename(pptx_path),
        "images": images,
        "n_images": len(images),
        "fonts_used": [
            {"family": f, "count": c} for f, c in font_families.most_common()
        ],
    }
    _write_yaml(outdir / "asset_manifest.yaml", manifest)

    # === 옵션 E: slides_full_dump.yaml ===
    if detail_full:
        full = {"source": os.path.basename(pptx_path), "slides": slides_dump}
        _write_yaml(outdir / "slides_full_dump.yaml", full)

    print(f"✅ 추출 완료 → {outdir}")
    print(f"   style_spec.yaml ({len(layout_counter)} layouts)")
    print(f"   layout_catalog.md")
    print(f"   narrative_outline.md")
    print(f"   asset_manifest.yaml ({len(images)} images)")
    if detail_full:
        print(f"   slides_full_dump.yaml ({len(slides_dump)} slides, --detail full)")


def _summarize_runs(runs: list[dict], top_n: int = 1) -> dict | list[dict]:
    """runs 리스트에서 가장 흔한 폰트 family/size를 요약."""
    if not runs:
        return {} if top_n == 1 else []
    fam_c: Counter = Counter()
    size_c: Counter = Counter()
    bold_c: Counter = Counter()
    color_c: Counter = Counter()
    for r in runs:
        f = r["font"]
        if f["family"]:
            fam_c[f["family"]] += 1
        if f["size_pt"]:
            size_c[round(f["size_pt"])] += 1
        if f["bold"] is not None:
            bold_c[f["bold"]] += 1
        if f["color_hex"]:
            color_c[f["color_hex"]] += 1

    def _top(c: Counter):
        return c.most_common(1)[0][0] if c else None

    summary = {
        "family": _top(fam_c),
        "size_pt": _top(size_c),
        "bold": _top(bold_c),
        "color_hex": _top(color_c),
    }
    return summary if top_n == 1 else [summary]


def _write_yaml(path: Path, obj) -> None:
    """yaml 미설치 환경 대비 — 간단한 자체 dump (set/dict/list/str/int/float/None만)."""
    try:
        import yaml  # type: ignore
        path.write_text(
            yaml.safe_dump(obj, sort_keys=False, allow_unicode=True),
            encoding="utf-8",
        )
        return
    except ImportError:
        pass

    # fallback: 자체 dump
    def _dump(o, indent=0):
        sp = "  " * indent
        if isinstance(o, dict):
            if not o:
                return "{}\n"
            out = []
            for k, v in o.items():
                if isinstance(v, (dict, list)) and v:
                    out.append(f"{sp}{k}:\n{_dump(v, indent+1)}")
                else:
                    out.append(f"{sp}{k}: {_scalar(v)}\n")
            return "".join(out)
        if isinstance(o, list):
            if not o:
                return f"{sp}[]\n"
            out = []
            for item in o:
                if isinstance(item, (dict, list)):
                    sub = _dump(item, indent+1)
                    out.append(f"{sp}-\n{sub}")
                else:
                    out.append(f"{sp}- {_scalar(item)}\n")
            return "".join(out)
        return f"{sp}{_scalar(o)}\n"

    def _scalar(v):
        if v is None:
            return "null"
        if isinstance(v, bool):
            return "true" if v else "false"
        if isinstance(v, (int, float)):
            return str(v)
        s = str(v)
        # 안전한 quote
        if any(ch in s for ch in [":", "#", "\n", "[", "]", "{", "}"]):
            s = s.replace('"', '\\"')
            return f'"{s}"'
        return s

    path.write_text(_dump(obj), encoding="utf-8")


def main() -> int:
    ap = argparse.ArgumentParser(description="Extract design spec from .pptx")
    ap.add_argument("pptx", help="입력 .pptx 경로")
    ap.add_argument("--out", default=None, help="출력 디렉토리 (기본: output/<deck-slug>/)")
    ap.add_argument("--detail", choices=["full"], default=None,
                    help="full: 슬라이드별 완전 dump 추가 (G3)")
    args = ap.parse_args()

    if not os.path.exists(args.pptx):
        sys.stderr.write(f"ERROR: 파일 없음 — {args.pptx}\n")
        return 1

    if args.out:
        outdir = Path(args.out)
    else:
        slug = Path(args.pptx).stem
        outdir = Path("output") / slug

    extract(args.pptx, outdir, detail_full=(args.detail == "full"))
    return 0


if __name__ == "__main__":
    sys.exit(main())
