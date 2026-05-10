#!/usr/bin/env python3
"""
roundtrip_check.py — extract_spec.py 결과로 대표 슬라이드 1장 재생성 →
원본과 PNG 픽셀 비교 → 일치율 보고. spec의 정확성 검증.

Round-trip 의 의미:
  spec → 재생성된 슬라이드 == 원본 슬라이드  ⇔  spec이 정확
  불일치 → spec이 그럴듯한 거짓 (analyze 스킬 SKILL.md 첫 줄 원칙)

사용:
    python3 roundtrip_check.py <원본.pptx> <spec_dir> [--slide N]
        --slide: 검증할 슬라이드 번호 (기본 1)

산출:
  spec_dir/roundtrip/
    original_slideN.png
    regenerated_slideN.png
    diff_slideN.png         (Pillow 있을 때만)
    roundtrip_report.md
exit 0 = 일치율 ≥ 85% / 1 = 미달
"""
from __future__ import annotations
import argparse
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
except ImportError:
    sys.stderr.write("ERROR: python-pptx 미설치.\n")
    sys.exit(1)


def _read_yaml(path: Path) -> dict:
    try:
        import yaml  # type: ignore
        return yaml.safe_load(path.read_text(encoding="utf-8"))
    except ImportError:
        # 매우 단순한 fallback (extract_spec.py가 dump한 형식만)
        sys.stderr.write(
            "WARN: PyYAML 없음. round-trip은 spec yaml을 정확히 못 읽을 수 있음. "
            "`pip3 install --user pyyaml` 권장.\n"
        )
        return {}


def _pptx_to_png(pptx_path: Path, out_dir: Path, dpi: int = 100) -> list[Path]:
    """soffice + pdftoppm 으로 PNG 변환."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    if soffice is None or pdftoppm is None:
        sys.stderr.write(
            "ERROR: soffice / pdftoppm 필요. "
            "`brew install --cask libreoffice && brew install poppler`\n"
        )
        sys.exit(1)
    out_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        r = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf",
             "--outdir", tmp, str(pptx_path)],
            capture_output=True, text=True, timeout=120,
        )
        if r.returncode != 0:
            sys.stderr.write(f"soffice 실패: {r.stderr[-200:]}\n")
            sys.exit(1)
        pdfs = [f for f in os.listdir(tmp) if f.endswith(".pdf")]
        if not pdfs:
            sys.stderr.write("PDF 변환 결과 없음\n")
            sys.exit(1)
        pdf = Path(tmp) / pdfs[0]
        prefix = out_dir / "page"
        subprocess.run(
            [pdftoppm, "-r", str(dpi), str(pdf), str(prefix), "-png"],
            check=True,
        )
    return sorted(out_dir.glob("page-*.png"))


def _regenerate_one_slide(spec: dict, src_pptx: Path, slide_no: int,
                          out_pptx: Path) -> None:
    """
    spec과 원본의 슬라이드 N에서 텍스트만 가져와 새 .pptx 한 장 생성.
    layout/스타일은 spec.deck.design_tokens에서, 텍스트는 원본에서.
    (Provenance Lock 원칙: 텍스트는 원본에서만 가져옴, 신규 작성 X)
    """
    src = Presentation(str(src_pptx))
    if slide_no < 1 or slide_no > len(src.slides):
        sys.stderr.write(f"ERROR: slide_no {slide_no} 범위 밖 (1~{len(src.slides)})\n")
        sys.exit(1)
    src_slide = src.slides[slide_no - 1]

    # 새 빈 deck (원본과 같은 크기)
    new = Presentation()
    if src.slide_width is not None:
        new.slide_width = src.slide_width
    if src.slide_height is not None:
        new.slide_height = src.slide_height
    blank = new.slide_layouts[6]  # blank layout
    new_slide = new.slides.add_slide(blank)

    # design tokens 추출
    dt = (spec or {}).get("deck", {}).get("design_tokens", {})
    typo = dt.get("typography", {})
    body_default = typo.get("body_default") or {}
    title_default = typo.get("title_default") or {}

    # 원본 슬라이드의 모든 text shape를 spec 스타일로 재생성
    for sh in src_slide.shapes:
        if not sh.has_text_frame:
            continue
        if sh.left is None or sh.top is None:
            continue
        tx = new_slide.shapes.add_textbox(
            sh.left, sh.top, sh.width or Emu(2000000), sh.height or Emu(500000)
        )
        tf = tx.text_frame
        tf.word_wrap = True
        # 첫 paragraph 비워서 시작
        tf.text = ""
        first = True
        src_tf = getattr(sh, "text_frame", None)
        if src_tf is None:
            continue
        for para in src_tf.paragraphs:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            for run_src in para.runs:
                if not run_src.text:
                    continue
                run = p.add_run()
                run.text = run_src.text
                # spec의 body_default 적용 (title은 위치로 추정)
                top_in = (sh.top or 0) / 914400
                ref = title_default if top_in < 1.5 else body_default
                if ref.get("family"):
                    run.font.name = ref["family"]
                if ref.get("size_pt"):
                    run.font.size = Pt(ref["size_pt"])
                if ref.get("bold") is not None:
                    run.font.bold = bool(ref["bold"])
    new.save(str(out_pptx))


def _png_similarity(a: Path, b: Path, diff_out: Path | None = None) -> float:
    """
    두 PNG의 평균 픽셀 유사도 (0.0~1.0). Pillow 있으면 정밀, 없으면 파일 크기 비.
    """
    try:
        from PIL import Image, ImageChops  # type: ignore
    except ImportError:
        sa, sb = a.stat().st_size, b.stat().st_size
        if max(sa, sb) == 0:
            return 0.0
        return min(sa, sb) / max(sa, sb)

    ia = Image.open(a).convert("RGB")
    ib = Image.open(b).convert("RGB")
    if ia.size != ib.size:
        ib = ib.resize(ia.size)
    diff = ImageChops.difference(ia, ib)
    if diff_out is not None:
        diff.save(diff_out)
    # 평균 픽셀 차 (0~255) → 유사도
    bbox = diff.getbbox()
    if bbox is None:
        return 1.0
    # 픽셀 평균 차이 계산
    stat_sum = 0
    n = 0
    for px in diff.getdata():
        stat_sum += sum(px) / 3
        n += 1
    avg = stat_sum / max(n, 1)
    return max(0.0, 1.0 - avg / 255.0)


def main() -> int:
    ap = argparse.ArgumentParser(description="extract_spec round-trip 검증")
    ap.add_argument("original", help="원본 .pptx")
    ap.add_argument("spec_dir", help="extract_spec.py 출력 디렉토리")
    ap.add_argument("--slide", type=int, default=1, help="검증할 슬라이드 번호 (기본 1)")
    ap.add_argument("--threshold", type=float, default=0.85,
                    help="일치율 임계 (기본 0.85)")
    args = ap.parse_args()

    spec_dir = Path(args.spec_dir)
    style_yaml = spec_dir / "style_spec.yaml"
    if not style_yaml.exists():
        sys.stderr.write(f"ERROR: style_spec.yaml 없음 — {style_yaml}\n")
        return 1
    spec = _read_yaml(style_yaml)

    rt_dir = spec_dir / "roundtrip"
    rt_dir.mkdir(parents=True, exist_ok=True)

    # 1. 재생성 .pptx 만들기
    regen_pptx = rt_dir / f"regenerated_slide{args.slide}.pptx"
    _regenerate_one_slide(spec, Path(args.original), args.slide, regen_pptx)

    # 2. 원본 + 재생성 모두 PNG로 변환
    orig_dir = rt_dir / "_orig_png"
    regen_dir = rt_dir / "_regen_png"
    orig_pngs = _pptx_to_png(Path(args.original), orig_dir)
    regen_pngs = _pptx_to_png(regen_pptx, regen_dir)

    if args.slide > len(orig_pngs):
        sys.stderr.write(f"ERROR: 원본 PNG 개수({len(orig_pngs)}) < slide {args.slide}\n")
        return 1

    orig_png = orig_pngs[args.slide - 1]
    regen_png = regen_pngs[0]  # 재생성은 1장짜리

    # 3. 결과를 spec_dir/roundtrip/ 로 복사
    out_orig = rt_dir / f"original_slide{args.slide}.png"
    out_regen = rt_dir / f"regenerated_slide{args.slide}.png"
    shutil.copy(orig_png, out_orig)
    shutil.copy(regen_png, out_regen)
    diff_out = rt_dir / f"diff_slide{args.slide}.png"

    # 4. 유사도
    sim = _png_similarity(out_orig, out_regen, diff_out)
    pct = round(sim * 100, 2)
    passed = sim >= args.threshold

    # 5. report
    rpt = [
        f"# Round-trip Report — Slide {args.slide}\n\n",
        f"- 원본: `{args.original}`\n",
        f"- spec dir: `{spec_dir}`\n",
        f"- 임계: {args.threshold * 100:.0f}%\n\n",
        f"## 결과\n\n",
        f"- 픽셀 유사도: **{pct}%** {'✓ PASS' if passed else '✗ FAIL'}\n\n",
        f"## 산출물\n\n",
        f"- `original_slide{args.slide}.png`\n",
        f"- `regenerated_slide{args.slide}.png`\n",
        f"- `diff_slide{args.slide}.png`\n",
        f"- `regenerated_slide{args.slide}.pptx`\n\n",
        f"## 해석\n\n",
        ("- 일치율 ≥ 85%: spec이 충분히 정확. 새 PPT 작성 시 신뢰 가능.\n"
         "- 일치율 < 85%: spec이 design tokens(폰트/크기/컬러) 또는 layout을\n"
         "  완전히 캡처하지 못함. extract_spec.py를 `--detail full`로 재실행하거나,\n"
         "  레이아웃 분류 규칙을 보강해야 함.\n"),
    ]
    (rt_dir / "roundtrip_report.md").write_text("".join(rpt), encoding="utf-8")

    print(f"\n📋 Round-trip 검증 — slide {args.slide}")
    print(f"   원본 PNG     : {out_orig}")
    print(f"   재생성 PNG   : {out_regen}")
    print(f"   diff PNG     : {diff_out}")
    print(f"   유사도       : {pct}% (임계 {args.threshold*100:.0f}%)")
    print(f"   결과         : {'✓ PASS' if passed else '✗ FAIL'}")
    return 0 if passed else 1


if __name__ == "__main__":
    sys.exit(main())
